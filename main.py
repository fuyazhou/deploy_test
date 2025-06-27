import os
from pptx import Presentation
from pptx.util import Inches, Pt
# Inches 和 Pt 可能在后续填充内容时用到，先导入

from langchain_openai import ChatOpenAI

from ppt_utils import load_ppt_layouts, generate_outline_from_llm, choose_layout_for_outline_item

# --- 配置区 ---
# 用户可以修改以下配置来自定义脚本行为
# ==============================================================================
# 1. PPT模板路径:
#    指定一个PPTX文件作为模板。脚本将使用此模板中的版式和设计。
#    如果文件不存在，脚本会尝试创建一个名为 "dummy_template.pptx" 的基础模板。
#    强烈建议使用您自己的、包含多种版式设计（如标题页、内容页、章节标题页等）的模板。
PPT_TEMPLATE_PATH = "my_custom_template.pptx"  # <--- 修改这里，指向你的模板文件

# 2. 输出PPT文件名:
#    生成的演示文稿将以此名称保存。
DEFAULT_OUTPUT_PPT_NAME = "generated_ai_presentation.pptx" # <--- 可修改输出文件名

# 3. 大语言模型 (LLM) 配置:
#    这些设置用于连接和控制AI模型的行为。
LLM_BASE_URL = "http://10.120.3.51:11436/v1"  # <--- 修改这里，如果你的LLM服务地址不同
LLM_MODEL_NAME = "Qwen/Qwen2.5-72B-Instruct-GPTQ-Int4" # <--- 可根据你的模型修改
LLM_TEMPERATURE = 0.1  # 控制输出的随机性，0表示更确定性，1表示更随机
# OPENAI_API_KEY = "sk-yourkeyhere" # 如果你的LLM服务需要API Key，请取消注释并填入

# 4. 生成大纲时的幻灯片数量建议:
#    这会提示LLM生成大约多少张幻灯片的大纲。实际数量可能略有不同。
MAX_SLIDES_FOR_OUTLINE = 7 # <--- 可调整期望的幻灯片数量

# 5. 详细模式:
#    设置为 True 会在控制台打印详细的执行步骤信息，方便调试。
#    设置为 False 则只打印关键信息和最终结果。
VERBOSE_MODE = True
# ==============================================================================


# 如果 PPT_TEMPLATE_PATH 指定的文件不存在，则使用此名称创建基础模板
DEFAULT_DUMMY_TEMPLATE_NAME = "dummy_template.pptx"
# 注意：DEFAULT_OUTPUT_PPT_NAME 已在上面定义，此处是重复的，将由上面的定义覆盖
# LLM_BASE_URL, LLM_MODEL_NAME 等也已在上面定义
LLM_TEMPERATURE = 0.1
# OPENAI_API_KEY = "sk-proj-..." # 如果需要，取消注释并填入

MAX_SLIDES_FOR_OUTLINE = 5 # 大纲生成时建议的幻灯片数量
VERBOSE_MODE = True # 是否打印详细过程信息

def v_print(message):
    """Verbose print function"""
    if VERBOSE_MODE:
        print(message)

def create_basic_template_if_not_exists(target_template_path, fallback_dummy_name):
    """
    如果 target_template_path 指定的模板文件不存在，
    则在当前目录创建一个名为 fallback_dummy_name 的非常基础的PPTX文件作为备用。
    """
    if not os.path.exists(target_template_path):
        v_print(f"User-defined template '{target_template_path}' not found.")
        v_print(f"Attempting to create a basic fallback template named '{fallback_dummy_name}'.")
        try:
            prs = Presentation()
            # python-pptx 默认会创建一组标准版式，如：
            # 0: Title Slide, 1: Title and Content, 2: Section Header, 3: Two Content,
            # 4: Comparison, 5: Title Only, 6: Blank, 7: Content with Caption, 8: Picture with Caption
            prs.save(fallback_dummy_name)
            v_print(f"Basic fallback template '{fallback_dummy_name}' created successfully.")
            v_print(f"The script will USE '{fallback_dummy_name}' for this run.")
            return fallback_dummy_name # 返回实际使用的模板路径
        except Exception as e:
            print(f"CRITICAL: Error creating basic fallback template '{fallback_dummy_name}': {e}")
            print("Please ensure you have write permissions or provide a valid template path.")
            return None # 表示创建失败
    return target_template_path # 用户指定的模板存在，直接使用它

def populate_slide_content(slide, slide_data):
    """
    用大纲中的数据填充幻灯片。
    这是一个基础实现，可能需要根据具体模板的占位符进行调整。
    """
    v_print(f"  Populating slide with title: '{slide_data.get('title', '')}'")

    # 尝试填充标题
    try:
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", " ") # PPTX要求非None
        else:
            # 如果没有专门的标题占位符，尝试寻找第一个最大的占位符作为标题
            # 这部分逻辑可以更复杂，例如基于 placeholder_format.type
            title_shape = None
            for shape in slide.placeholders:
                if shape.is_placeholder and shape.name.lower().startswith(("title", "标题")):
                    title_shape = shape
                    break
            if title_shape:
                 title_shape.text_frame.text = slide_data.get("title", " ")
                 v_print(f"    Filled title placeholder '{title_shape.name}' with '{slide_data.get('title', '')}'")
            elif slide.placeholders: # 拿第一个可用的占位符作为标题（不太理想但作为后备）
                slide.placeholders[0].text_frame.text = slide_data.get("title", " ")
                v_print(f"    Used first placeholder '{slide.placeholders[0].name}' for title.")


    except Exception as e:
        v_print(f"    Warning: Could not set title for slide. Error: {e}")

    # 尝试填充内容要点
    points = slide_data.get("points", [])
    if points:
        body_shape = None
        # 寻找内容占位符 (通常名为 "Body Placeholder", "Content Placeholder", "内容占位符" 等)
        # 或者根据索引，通常标题后的第一个大占位符是内容区
        for shape in slide.placeholders:
            # Placeholder Type: BODY (1), CONTENT (1), TEXT_BOX (14) etc.
            # We are looking for a placeholder that is not the title and is suitable for body text.
            if shape.is_placeholder and shape.name and \
               ("body" in shape.name.lower() or \
                "content" in shape.name.lower() or \
                "内容" in shape.name.lower() or \
                (shape.placeholder_format.type == 1 and shape != slide.shapes.title) or # BODY type
                (shape.placeholder_format.type == 7 and shape != slide.shapes.title) # CONTENT type, often includes more than just text
                ):
                # 检查是否是标题占位符，如果 slide.shapes.title 存在且与当前 shape 相同，则跳过
                if slide.shapes.title and shape == slide.shapes.title:
                    continue
                body_shape = shape
                break

        if not body_shape and len(slide.placeholders) > 1:
            # 如果没有明确找到，并且有多个占位符，尝试使用非标题的第一个占位符
            # 这假设第一个占位符是标题（如果存在），第二个是内容
            potential_body_shapes = [p for p in slide.placeholders if p != slide.shapes.title]
            if potential_body_shapes:
                body_shape = potential_body_shapes[0]

        if body_shape:
            tf = body_shape.text_frame
            tf.clear() # 清除占位符中的默认文本

            # 第一个point可能作为段落标题或直接作为第一点
            # p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
            # p.text = points[0]
            # p.font.bold = True # 可选：将第一点加粗

            # for point_text in points: # (如果第一点已处理，则从 points[1:] 开始)
            #     p = tf.add_paragraph()
            #     p.text = point_text
            #     p.level = 0 # 默认级别，可以根据需要设置缩进级别

            # 更简洁的填充方式：
            tf.text = "\n".join(points) # 将所有点用换行符连接起来直接设置
            v_print(f"    Filled body placeholder '{body_shape.name}' with {len(points)} points.")
            # 你可能需要更细致地处理段落格式，比如字号、缩进等
            for paragraph in tf.paragraphs:
                 if not paragraph.runs: #确保有run
                     paragraph.add_run()
                 paragraph.font.size = Pt(18) # 示例字号
        else:
            v_print("    Warning: Could not find a suitable body placeholder for points on this slide layout.")
            # 如果没有找到合适的占位符，可以考虑添加一个新的文本框，但这会脱离模板的版式
            # left = top = width = height = Inches(1.0)
            # txBox = slide.shapes.add_textbox(left, top, width, height)
            # tf = txBox.text_frame
            # tf.text = "\n".join(points)


def main():
    print("=============================================")
    print("🤖 AI PPT Generator by Jules 🤖")
    print("=============================================")
    v_print("Starting the AI PPT Generation process...")

    # 0. 检查/创建模板文件
    #    create_basic_template_if_not_exists 会返回实际使用的模板路径
    #    如果用户指定的PPT_TEMPLATE_PATH不存在，它会尝试创建DEFAULT_DUMMY_TEMPLATE_NAME
    #    并返回DEFAULT_DUMMY_TEMPLATE_NAME作为实际使用的路径。
    v_print("\nStep 0: Verifying template file...")
    actual_template_path = create_basic_template_if_not_exists(PPT_TEMPLATE_PATH, DEFAULT_DUMMY_TEMPLATE_NAME)
    if not actual_template_path:
        print(f"CRITICAL: Failed to find or create a template. Exiting.")
        print(f"Please check the path '{PPT_TEMPLATE_PATH}' or permissions to create '{DEFAULT_DUMMY_TEMPLATE_NAME}'.")
        return

    v_print(f"Using template: '{actual_template_path}'")

    # 1. 初始化LLM
    v_print("\nStep 1: Initializing Language Model...")
    try:
        llm = ChatOpenAI(
            temperature=LLM_TEMPERATURE,
            model=LLM_MODEL_NAME,
            base_url=LLM_BASE_URL,
            # openai_api_key=OPENAI_API_KEY # 如果需要
        )
        v_print("Language Model initialized successfully.")
    except Exception as e:
        print(f"Error initializing Language Model: {e}")
        return

    # 2. 加载PPT模板版式
    v_print(f"\nStep 2: Loading PPT layouts from '{actual_template_path}'...")
    layouts = load_ppt_layouts(actual_template_path)
    if not layouts:
        print(f"No layouts loaded from '{actual_template_path}'. This could be due to an invalid or empty PPTX file.")
        print("If you are using the dummy template, it should contain default layouts.")
        print("Please check the template file or an AGENTS.md file for more info. Exiting.")
        return
    v_print(f"Loaded {len(layouts)} layouts: {', '.join(list(layouts.keys()))}")

    # 3. 获取用户主题并生成大纲
    print("\n---------------------------------------------")
    user_topic = input("Enter the topic for your presentation: ")
    if not user_topic:
        print("No topic entered. Exiting.")
        return

    v_print(f"\nStep 3: Generating outline for topic: '{user_topic}' (max_slides: {MAX_SLIDES_FOR_OUTLINE})...")
    outline = generate_outline_from_llm(user_topic, llm, max_slides=MAX_SLIDES_FOR_OUTLINE)

    if not outline:
        print("Failed to generate outline. Exiting.")
        return
    v_print("Outline generated successfully:")
    for i, item in enumerate(outline):
        v_print(f"  Slide {i+1}: Title: {item.get('title', 'N/A')}, Points: {len(item.get('points', []))}")

    # 4. 创建新的演示文稿 (可以基于模板)
    v_print(f"\nStep 4: Creating new presentation (based on '{actual_template_path}')...")
    try:
        # 使用实际的模板路径（可能是用户指定的，也可能是创建的dummy模板）
        prs = Presentation(actual_template_path)

        # 关于处理模板中现有幻灯片的说明:
        # 当前实现会直接在模板的现有幻灯片之后添加新幻灯片。
        # 如果希望每次都从一个“干净”的演示开始（仅使用模板的母版和设计，不含其内容幻灯片），
        # 你有几种选择：
        # 1. 最佳实践: 准备一个“干净”的模板文件，该文件只包含设计好的母版和版式，不包含任何实际的幻灯片页。
        # 2. 代码中删除: 可以取消下面代码块的注释，以在运行时尝试删除模板中的所有现有幻灯片。
        #    这有一定风险，如果模板结构特殊或为空，可能会出错。
        #
        # if len(prs.slides._sldIdLst) > 0: # 检查是否有幻灯片
        #     v_print(f"Template '{actual_template_path}' has {len(prs.slides)} existing slides.")
        #     v_print("Attempting to clear existing slides to start fresh (this is experimental)...")
        #     try:
        #         for i in range(len(prs.slides) - 1, -1, -1): # 从后往前删
        #             rId = prs.slides._sldIdLst[i].rId
        #             prs.part.drop_rel(rId)
        #             del prs.slides._sldIdLst[i]
        #         v_print("Successfully cleared existing slides.")
        #     except Exception as clear_e:
        #         v_print(f"Could not clear all existing slides: {clear_e}. New slides will be appended.")
        # else:
        #     v_print(f"Template '{actual_template_path}' is clean (no existing slides).")

    except Exception as e:
        v_print(f"Error opening presentation with template '{actual_template_path}': {e}")
        v_print("Creating a new presentation from scratch instead (will lack custom template styling).")
        prs = Presentation() # Fallback to a truly blank presentation

    v_print("New presentation object created.")

    # 5. 遍历大纲，选择版式，添加幻灯片并填充内容
    v_print("\nStep 5: Processing outline and generating slides...")
    for i, slide_data in enumerate(outline):
        v_print(f"\nProcessing Slide {i+1}/{len(outline)}: '{slide_data.get('title', 'Untitled')}'")

        # a. 选择版式
        chosen_layout_name = choose_layout_for_outline_item(slide_data, layouts, llm)

        selected_layout = None
        if chosen_layout_name and chosen_layout_name in layouts:
            selected_layout = layouts[chosen_layout_name]
            v_print(f"  Layout chosen by LLM: '{chosen_layout_name}'")
        else:
            v_print(f"  LLM did not choose a valid layout ('{chosen_layout_name}'). Trying a default.")
            # 尝试使用一个常见的默认版式，例如“Title and Content”
            # 注意：这里的名称需要与你的模板中的版式名称匹配
            default_layout_names_to_try = ["Title and Content", "标题和内容", layouts[1].name if len(layouts)>1 else list(layouts.keys())[0]]
            for name in default_layout_names_to_try:
                if name in layouts:
                    selected_layout = layouts[name]
                    v_print(f"  Using default layout: '{name}'")
                    break
            if not selected_layout: # 如果连默认的都找不到，就用第一个可用的版式
                 selected_layout = layouts[list(layouts.keys())[0]]
                 v_print(f"  Using first available layout as fallback: '{list(layouts.keys())[0]}'")


        # b. 添加幻灯片
        try:
            slide_added = prs.slides.add_slide(selected_layout)
            v_print(f"  Added slide with layout: '{selected_layout.name if selected_layout else 'N/A'}'")
        except Exception as e:
            print(f"  Error adding slide with layout '{selected_layout.name if selected_layout else 'N/A'}': {e}")
            print(f"  Attempting to add slide with the first available layout from the base presentation.")
            try:
                # Fallback: use a layout from the presentation's own slide master if template layout fails
                # This usually corresponds to the default layouts of a blank presentation.
                fallback_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
                slide_added = prs.slides.add_slide(fallback_layout)
                v_print(f"  Added slide with presentation's internal layout: '{fallback_layout.name}'")
            except Exception as e2:
                print(f"  FATAL: Could not add slide even with internal fallback layout: {e2}")
                continue # Skip this slide


        # c. 填充内容
        populate_slide_content(slide_added, slide_data)

    # 6. 保存PPT
    output_filename = DEFAULT_OUTPUT_PPT_NAME
    try:
        prs.save(output_filename)
        v_print(f"\nStep 6: Presentation saved as '{output_filename}'")
        print(f"\n✅ Presentation '{output_filename}' generated successfully!")
    except Exception as e:
        print(f"Error saving presentation: {e}")

if __name__ == '__main__':
    main()
