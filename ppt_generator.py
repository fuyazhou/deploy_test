import json
import re
from typing import Dict, List

from pptx import Presentation

try:
    from langchain_openai import ChatOpenAI
except ImportError:
    ChatOpenAI = None


def parse_template(path: str) -> Dict[str, Dict[str, List[str]]]:
    """Parse slide layouts from a PowerPoint template.

    Returns a mapping of layout name to a dictionary describing
    the placeholders available in that layout.
    """
    prs = Presentation(path)
    layouts = {}
    for i, layout in enumerate(prs.slide_layouts):
        placeholders = [ph.name for ph in layout.placeholders]
        layouts[f"layout_{i}"] = {
            "placeholders": placeholders,
        }
    return layouts


def call_llm(prompt: str) -> str:
    """Call the language model with the provided prompt."""
    if ChatOpenAI is None:
        raise RuntimeError("langchain_openai is not installed")
    qwen = ChatOpenAI(
        temperature=0,
        model="Qwen/Qwen2.5-72B-Instruct-GPTQ-Int4",
        base_url="http://10.120.3.51:11436/v1",
    )
    return qwen.invoke(prompt).content


def parse_json_from_text(text: str):
    """Extract the first JSON object found in text."""
    match = re.search(r"\{.*\}", text, re.S)
    if match:
        return json.loads(match.group())
    raise ValueError("No JSON object found in response")


def choose_layout(outline_item: str, layouts: Dict[str, dict]) -> str:
    prompt = (
        "Given the outline item: '" + outline_item + "',\n"
        "Choose the most suitable layout from the following options:"\
        f" {list(layouts.keys())}.\n"\
        "Respond with the layout name in JSON like {\"layout\": \"layout_1\"}."
    )
    response = call_llm(prompt)
    data = parse_json_from_text(response)
    layout_name = data.get("layout")
    if layout_name not in layouts:
        layout_name = list(layouts.keys())[0]
    return layout_name


def create_presentation(
    outline: List[str],
    layouts: Dict[str, dict],
    output_path: str,
    template_path: str,
):
    prs = Presentation(template_path)
    for item in outline:
        layout_name = choose_layout(item, layouts)
        layout_index = int(layout_name.split("_")[-1])
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = item
    prs.save(output_path)


def main(user_prompt: str, template_path: str, output_path: str):
    layouts = parse_template(template_path)
    outline_prompt = (
        "Create a JSON list describing the slides for the following topic:\n"
        + user_prompt
    )
    response = call_llm(outline_prompt)
    data = parse_json_from_text(response)
    outline = data.get("slides", [])
    if not isinstance(outline, list):
        outline = []
    create_presentation(outline, layouts, output_path, template_path)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Generate a PowerPoint from a prompt.")
    parser.add_argument("prompt", help="User prompt for the presentation")
    parser.add_argument("template", help="Path to PowerPoint template")
    parser.add_argument("output", help="Path for the generated presentation")
    args = parser.parse_args()

    main(args.prompt, args.template, args.output)
