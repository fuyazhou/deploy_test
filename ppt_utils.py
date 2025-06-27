import json
from pptx import Presentation
# 假设 ChatOpenAI 会从 langchain_openai 导入，如果用户环境不同，可能需要调整
# from langchain_openai import ChatOpenAI

# 为了模块的独立性，ChatOpenAI 实例将作为参数传入，而不是在此处直接初始化。
# 这也使得函数更易于测试。

def load_ppt_layouts(template_path: str) -> dict:
    """
    加载PPT模板文件中的所有幻灯片版式。

    参数:
        template_path (str): PPT模板文件的路径。

    返回:
        dict: 一个字典，其中键是版式名称 (str)，值是版式对象 (SlideLayout)。
              如果文件未找到或发生错误，则返回空字典。
    """
    layouts = {}
    try:
        prs = Presentation(template_path)
        for i, layout in enumerate(prs.slide_layouts):
            layout_name = layout.name
            # 如果版式名称为空或者重复，则生成一个唯一的名称
            if not layout_name:
                layout_name = f"UnnamedLayout_{i+1}"
            elif layout_name in layouts:
                original_name = layout_name
                count = 1
                while layout_name in layouts:
                    layout_name = f"{original_name}_{count}"
                    count += 1
            layouts[layout_name] = layout
        return layouts
    except Exception as e:
        print(f"Error loading PPT template '{template_path}': {e}")
        return {}

def generate_outline_from_llm(user_topic: str, llm_instance, max_slides: int = 7) -> list | None:
    """
    使用大模型根据用户主题生成PPT大纲。

    参数:
        user_topic (str): 用户提供PPT主题。
        llm_instance: 已初始化的ChatOpenAI语言模型实例。
        max_slides (int): 期望生成的最大幻灯片数量（大致数量）。

    返回:
        list: 解析后的大纲列表，每个元素是一个字典，代表一张幻灯片。
              例如: [{'title': 'Slide 1 Title', 'points': ['Point 1', 'Point 2']}, ...]
              如果发生错误或无法解析JSON，则返回None。
    """
    prompt = f"""请为关于主题“{user_topic}”的演示文稿创建一个幻灯片大纲。
我希望大纲是JSON格式的。JSON应该是一个列表，每个列表项代表一张幻灯片。
每张幻灯片应该有一个'title'键（字符串类型）和一个'points'键（字符串列表类型，包含该幻灯片的主要内容点）。
请确保生成的JSON是有效的，并且可以直接被Python的json.loads()解析。
请尽量生成大约 {max_slides} 张幻灯片的大纲。

JSON输出示例:
[
  {{
    "title": "幻灯片1：导论",
    "points": ["介绍主题", "概述目标"]
  }},
  {{
    "title": "幻灯片2：主要概念",
    "points": ["解释概念A", "详细说明概念B", "概念A与概念B的关系"]
  }}
]

请严格按照此JSON格式提供大纲:
"""

    try:
        response = llm_instance.invoke(prompt)
        content = response.content  # 或者 response.text，取决于具体实现

        # 尝试从模型输出中提取JSON块
        # 大模型有时会在JSON前后添加额外的文本，如 "Here is the JSON output:" 或 ```json ... ```
        json_start_index = content.find('[')
        json_end_index = content.rfind(']')

        if json_start_index != -1 and json_end_index != -1 and json_start_index < json_end_index:
            json_str = content[json_start_index : json_end_index+1]

            # 额外的清理，去除可能的markdown代码块标记
            if json_str.startswith("```json"):
                json_str = json_str[len("```json"):].strip()
            if json_str.endswith("```"):
                json_str = json_str[:-len("```")].strip()

            outline = json.loads(json_str)
            # 基本的验证，确保是大纲结构的列表，且列表内是字典
            if isinstance(outline, list) and all(isinstance(item, dict) for item in outline):
                return outline
            else:
                print(f"Parsed JSON is not in the expected format (list of dicts): {outline}")
                return None
        else:
            print(f"Could not find a valid JSON list in the LLM response: {content}")
            return None

    except json.JSONDecodeError as e:
        print(f"Error decoding JSON from LLM response: {e}")
        print(f"LLM raw response was: {content}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred while generating outline: {e}")
        print(f"LLM raw response was: {content if 'content' in locals() else 'N/A'}")
        return None

def choose_layout_for_outline_item(outline_item: dict, available_layouts: dict, llm_instance) -> str | None:
    """
    根据大纲项目的内容，使用大模型从可用版式中选择一个最合适的版式名称。

    参数:
        outline_item (dict): 单个大纲项目，期望包含 'title' 和可选的 'points' 键。
                             例如: {'title': '引言', 'points': ['点1', '点2']}
        available_layouts (dict): 包含可用版式名称的字典 (键是版式名称, 值是SlideLayout对象).
        llm_instance: 已初始化的ChatOpenAI语言模型实例。

    返回:
        str: 被选中的最合适的版式名称。如果无法做出选择或出错，则返回None。
             也可能返回一个默认的版式名称，如 "Title and Content"，如果LLM未能成功选择。
    """
    item_title = outline_item.get("title", "")
    item_points = outline_item.get("points", [])

    # 构建对大纲项目内容的简要描述
    content_summary = f"幻灯片标题: '{item_title}'."
    if item_points:
        content_summary += " 主要内容点包括: " + ", ".join(item_points[:3]) # 取前3个点作为摘要
        if len(item_points) > 3:
            content_summary += " 等。"
    else:
        content_summary += " (这张幻灯片可能没有详细的文本点，或者主要是标题/图片/图表)。"

    layout_names = list(available_layouts.keys())
    if not layout_names:
        print("Error: No available layouts provided to choose_layout_for_outline_item.")
        return None # 或者返回一个预设的默认值

    prompt = f"""我正在为一个演示文稿的特定幻灯片选择最佳版式。
幻灯片的内容如下：
{content_summary}

可用的幻灯片版式有：
{', '.join(layout_names)}

请根据幻灯片的内容，从以上列表中选择一个最合适的版式名称。
请只返回你选择的版式名称，不要添加任何额外的解释或文字。
例如，如果列表是 ["Title Slide", "Title and Content", "Blank"]，而你认为 "Title and Content" 最合适，则只返回 "Title and Content"。
你的选择："""

    try:
        response = llm_instance.invoke(prompt)
        chosen_layout_name = response.content.strip()

        # 清理LLM可能返回的额外引号或标点
        chosen_layout_name = chosen_layout_name.replace("\"", "").replace("'", "").replace(".", "")

        # 验证返回的名称是否在可用列表中
        if chosen_layout_name in layout_names:
            return chosen_layout_name
        else:
            # 尝试进行模糊匹配或选择一个默认值
            print(f"LLM chose a layout ('{chosen_layout_name}') not in the available list: {layout_names}.")
            # 简单的回退逻辑：如果LLM的选择无效，可以尝试选择一个常见的默认版式（如果存在）
            # 或者可以尝试找到最相似的名称，但这会更复杂
            default_candidates = ["Title and Content", "标题和内容", "内容与标题"] # 常见版式名称
            for candidate in default_candidates:
                if candidate in layout_names:
                    print(f"Falling back to default layout: '{candidate}'")
                    return candidate
            # 如果连默认的都不在，就返回列表中的第一个作为最后的手段
            if layout_names:
                print(f"Falling back to the first available layout: '{layout_names[0]}'")
                return layout_names[0]
            return None # 如果列表为空，则没办法了

    except Exception as e:
        print(f"An error occurred while choosing layout: {e}")
        # 同样，可以在出错时返回一个默认版式
        default_candidates = ["Title and Content", "标题和内容", "内容与标题"]
        for candidate in default_candidates:
            if candidate in layout_names:
                print(f"Error occurred, falling back to default layout: '{candidate}'")
                return candidate
        if layout_names: # 确保 layout_names 不为空
             print(f"Error occurred, falling back to the first available layout: '{layout_names[0]}'")
             return layout_names[0]
        return None


if __name__ == '__main__':
    # --- 测试 load_ppt_layouts ---
    try:
        prs_dummy = Presentation()
        prs_dummy.save("dummy_template.pptx")
        print("Created dummy_template.pptx for testing load_ppt_layouts.")
        layouts_dict = load_ppt_layouts("dummy_template.pptx")
        if layouts_dict:
            print("\nAvailable layouts from dummy_template.pptx:")
            for name in layouts_dict.keys():
                print(f"- '{name}'")
        else:
            print("No layouts found in dummy_template.pptx or error in loading.")
    except Exception as e:
        print(f"Error in dummy template creation or load_ppt_layouts test: {e}")

    # --- 测试 generate_outline_from_llm 和 choose_layout_for_outline_item ---
    qwen_llm_instance_for_test = None
    generated_outline_for_test = None

    # Setup LLM instance (common for multiple tests)
    # 注意: 这部分测试需要配置并能够访问到ChatOpenAI模型实例 (qwen)
    # 你需要取消注释并正确配置下面的 ChatOpenAI 初始化代码
    print("\n--- Setting up LLM instance for tests ---")
    try:
        from langchain_openai import ChatOpenAI
        qwen_llm_instance_for_test = ChatOpenAI(
            temperature=0.1,
            model="Qwen/Qwen2.5-72B-Instruct-GPTQ-Int4",
            base_url="http://10.120.3.51:11436/v1",
            # openai_api_key="YOUR_API_KEY_IF_NEEDED"
        )
        print("Qwen LLM instance created (or attempted).")
    except ImportError:
        print("langchain_openai or ChatOpenAI could not be imported. LLM-dependent tests will be skipped.")
    except Exception as e:
        print(f"Error during LLM setup: {e}. LLM-dependent tests will be affected.")
        print("Ensure your LLM (Qwen) is running and accessible at the specified base_url.")

    # Test generate_outline_from_llm
    print("\n--- Testing generate_outline_from_llm ---")
    if qwen_llm_instance_for_test:
        # test_topic = "人工智能在教育领域的应用"
        # print(f"Generating outline for topic: '{test_topic}' (LLM call will be skipped)...")
        # generated_outline_for_test = generate_outline_from_llm(test_topic, qwen_llm_instance_for_test, max_slides=3)
        # if generated_outline_for_test:
        #     print("Generated Outline (simulated):")
        #     for i, slide_data in enumerate(generated_outline_for_test):
        #         print(f"  Slide {i+1}: Title: {slide_data.get('title', 'N/A')}, Points: {slide_data.get('points', [])}")
        # else:
        #     print("Failed to generate outline (simulated).")
        print("Skipping actual LLM call for generate_outline_from_llm in this automated test environment.")
        # Use a manual dummy outline for subsequent tests if LLM call is skipped or fails
        generated_outline_for_test = [
            {"title": "引言：AI医疗", "points": ["医疗行业现状", "AI能解决的问题"]},
            {"title": "AI具体应用", "points": ["智能诊断", "药物研发", "个性化治疗"]},
            {"title": "仅含标题的幻灯片"}, # 测试没有 points 的情况
            {"title": "未来展望与挑战", "points": ["技术趋势", "伦理问题", "数据隐私", "更多要点1", "更多要点2"]}
        ]
        print(f"Using manual outline for subsequent tests: {json.dumps(generated_outline_for_test, ensure_ascii=False, indent=2)}")
    else:
        print("Skipping generate_outline_from_llm test as LLM instance is not available.")
        # Provide a fallback dummy outline if LLM instance failed
        generated_outline_for_test = [
            {"title": "引言 (no LLM)", "points": ["内容1"]},
            {"title": "结论 (no LLM)", "points": ["内容2"]}
        ]
        print(f"Using fallback manual outline: {json.dumps(generated_outline_for_test, ensure_ascii=False, indent=2)}")

    # Test choose_layout_for_outline_item
    print("\n--- Testing choose_layout_for_outline_item ---")
    if qwen_llm_instance_for_test and layouts_dict and generated_outline_for_test:
        print(f"Available layouts for choosing: {list(layouts_dict.keys())}")
        # for item in generated_outline_for_test:
        #     print(f"\nChoosing layout for outline item: Title='{item.get('title', 'N/A')}' (LLM call will be skipped)...")
        #     chosen_name = choose_layout_for_outline_item(item, layouts_dict, qwen_llm_instance_for_test)
        #     if chosen_name:
        #         print(f"  LLM (simulated) chose: '{chosen_name}'")
        #         if chosen_name in layouts_dict:
        #             print(f"  '{chosen_name}' is a valid layout.")
        #         else:
        #             print(f"  Warning: '{chosen_name}' is NOT in the list of available layouts (might be due to fallback).")
        #     else:
        #         print("  LLM (simulated) failed to choose a layout or an error occurred.")
        print("Skipping actual LLM call for choose_layout_for_outline_item in this automated test environment.")
        # Simulate some choices for demonstration
        if generated_outline_for_test and layouts_dict and list(layouts_dict.keys()):
            first_layout_name = list(layouts_dict.keys())[0]
            print(f"Simulated choice for '{generated_outline_for_test[0]['title']}': '{first_layout_name}'")
            if len(generated_outline_for_test) > 1:
                # Try to pick a different layout for the second item if available
                second_layout_name_idx = 1 if len(list(layouts_dict.keys())) > 1 else 0
                second_layout_name = list(layouts_dict.keys())[second_layout_name_idx]
                print(f"Simulated choice for '{generated_outline_for_test[1]['title']}': '{second_layout_name}'")
    elif not layouts_dict:
        print("Skipping choose_layout_for_outline_item test because layouts_dict is empty or not loaded.")
    elif not generated_outline_for_test:
        print("Skipping choose_layout_for_outline_item test because generated_outline_for_test is empty or was not created.")
    else: # This means qwen_llm_instance_for_test is None
        print("Skipping choose_layout_for_outline_item test as LLM instance is not available.")

    # Example of using a loaded layout object (from previous test code, kept for reference)
    # prs_new = Presentation()
    # if layouts_dict and "Title Slide" in layouts_dict:
    #     slide_layout_obj = layouts_dict["Title Slide"]
    #     slide = prs_new.slides.add_slide(slide_layout_obj)
    #     title_shape = slide.shapes.title
    #     if title_shape: title_shape.text = "Hello from loaded layout"
    #     prs_new.save("test_slide_from_layout.pptx")
    #     print("\nCreated test_slide_from_layout.pptx using a loaded layout (example code).")
